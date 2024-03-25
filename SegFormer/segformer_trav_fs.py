"""
https://huggingface.co/nvidia/segformer-b3-finetuned-cityscapes-1024-1024
https://huggingface.co/blog/fine-tune-segformer
torchrun --standalone --nproc_per_node=gpu SegFormer/segformer_trav_ddp.py
Few-shot training of segformer on trav dataset.
Freeze segformer, only finetune the fusion transformer model
Change data_root:
/home/qiyuan/2023spring/segmentation_indoor_images
/mnt/hdd/segmentation_indoor_images

"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import torch
import json
from huggingface_hub import hf_hub_download
import argparse
import numpy as np
import matplotlib.pyplot as plt
from transformers import SegformerForSemanticSegmentation
import multiprocessing as mp
from tqdm import tqdm
from torch.utils.data import DataLoader
from torch.nn import functional as F
import torch.nn as nn
from torch.optim import SGD
from torch.nn.parallel import DistributedDataParallel as DDP
from torch.utils.data import DistributedSampler, RandomSampler
from torch.distributed import init_process_group, destroy_process_group
import wandb

from local_datasets.cityscapes import trav_train_loader, trav_val_loader
from models.transformer import MultiHeadAttentionOne
from models.infonce import InfoNCE
from utils.schedulers import get_scheduler, create_lr_scheduler
from utils.optimizers import get_optimizer
# from utils.utils import fix_seeds, setup_cudnn, cleanup_ddp, setup_ddp
import utils.distributed_utils as utils

# os.environ['OMP_NUM_THREADS'] = f"{int(os.cpu_count() // torch.cuda.device_count())}"  # no, slower

def ddp_setup():
    init_process_group(backend="nccl")
    torch.cuda.set_device(int(os.environ["LOCAL_RANK"]))


def get_argparser():
    parser = argparse.ArgumentParser('Pytorch SegFormer Models training and evaluation script', add_help=False)

    # Datset Options
    parser.add_argument("--data_root", type=str, default='/mnt/hdd/segmentation_indoor_images',help="path to Dataset")
    parser.add_argument("--scenes", type=list, default=['elb', 'erb', 'uc', 'wh'],
                        choices=['elb', 'erb', 'uc', 'nh', 'woh'], help='Name of dataset')
    parser.add_argument("--image_size", type=list, default=480, help="input size, [480, 640]")
    parser.add_argument("--ignore_label", type=int, default=255, help="path to Dataset")
    parser.add_argument("--dataset", type=str, default='trav',choices=['cityscapes', 'trav'], help='Name of dataset')
    parser.add_argument("--num_classes", type=int, default=2, help="num classes (default: None)")
    parser.add_argument("--pin_mem", type=bool, default=True, help="Dataloader ping_memory")
    parser.add_argument("--batch_size", type=int, default=1,help='batch size (zeus:10, poseidon:4)')
    parser.add_argument("--val_batch_size", type=int, default=4,help='batch size for validation')
    parser.add_argument("--scale_min", type=float, default=0.5)
    parser.add_argument("--scale_max", type=float, default=2.0)
    parser.add_argument("--rot_min", type=float, default=-10)
    parser.add_argument("--rot_max", type=float, default=10)
    parser.add_argument("--mean", type=list, default=[0.5174, 0.4857, 0.5054])
    parser.add_argument("--std", type=list, default=[0.2726, 0.2778, 0.2861])
    parser.add_argument("--augmentations", type=list, default=['hor_flip', 'vert_flip', 'resize'])
    parser.add_argument("--train_split", type=int, default=0)
    parser.add_argument("--test_split", type=str, default='default')
    parser.add_argument("--random_shot", type=bool, default=False)
    parser.add_argument("--shot", type=int, default=1)

    # Model Options
    # parser.add_argument("--model", type=str, default='make_SegFormerB1', help='model name')
    parser.add_argument("--bottleneck_dim", type=int, default=768)
    parser.add_argument("--heads", type=int, default=1)
    parser.add_argument("--dropout", type=float, default=0.5)

    # Train Options
    parser.add_argument("--amp", type=bool, default=False, help='auto mixture precision, do not use') # There may be some problems when loading weights, such as: ComplexFloat
    parser.add_argument("--epochs", type=int, default=20, help='total training epochs')
    parser.add_argument("--device", type=str, default='cuda:0', help='device (cuda:0 or cpu)')
    parser.add_argument("--workers", type=int, default=4,
                        help='workers, set it equal 0 when run programs in win platform')
    parser.add_argument("--DDP", type=bool, default=False)
    parser.add_argument("--train_print_freq", type=int, default=50)
    parser.add_argument("--val_print_freq", type=int, default=50)
    parser.add_argument("--adapt_iter", type=int, default=200)

    # Loss Options
    parser.add_argument("--loss_fn_name", type=str, default='CrossEntropy', choices=['OhemCrossEntropy', 'CrossEntropy'])

    # Optimizer & LR-scheduler Options
    parser.add_argument("--optimizer", type=str, default='adamw')
    parser.add_argument("--lr", type=float, default=2.5e-4,
                        help="learning rate (default: 1e-3)")
    parser.add_argument("--weight_decay", type=float, default=1e-4,
                        help='weight decay (default: 1e-5)')

    parser.add_argument("--lr_scheduler", type=str, default='WarmupPolyLR')
    parser.add_argument("--lr_power", type=float, default=0.9)
    parser.add_argument("--lr_warmup", type=int, default=10)
    parser.add_argument("--lr_warmup_ratio", type=float, default=0.1)

    # save checkpoints
    parser.add_argument("--save_weights_dir", default='./save_weights', type=str,
                        help="restore from checkpoint")
    parser.add_argument("--resume", type=bool, default=False)
    parser.add_argument("--gpu_id", type=int, default=0, help="GPU ID")
    parser.add_argument("--save_dir", type=str, default='./', help='SummaryWriter save dir')
    parser.add_argument("--eval_interval", type=int, default=2, help="evaluation interval")
    parser.add_argument("--save_every", type=int, default=200, help="save interval")
    parser.add_argument("--load_pretrained", type=bool, default=True)
    parser.add_argument("--pretrained_path", type=str, default='nvidia/segformer-b3-finetuned-cityscapes-1024-1024')
    parser.add_argument("--snapshot_path", type=str, default='save_weights/trav_fs_cpc')
    return parser

class Trainer:
    def __init__(self, args) -> None:
        self.args = args
        # init train objects
        if args.DDP:
            self.gpu_id = int(os.environ["LOCAL_RANK"])
        else:
            self.gpu_id = args.device

        repo_id = "huggingface/label-files"
        id2label_filename = "cityscapes-id2label.json"
        id2label = json.load(open(hf_hub_download(repo_id, id2label_filename, repo_type="dataset"), "r"))
        id2label = {int(k): v for k, v in id2label.items()}
        label2id = {v: k for k, v in id2label.items()}

        self.trainloader, _ = trav_train_loader(args)
        self.valloader, _ = trav_val_loader(args)

        self.model = SegformerForSemanticSegmentation.from_pretrained(args.pretrained_path,
            id2label=id2label, label2id=label2id)
        new_classifier = nn.Conv2d(768, 2, kernel_size=(1, 1), stride=(1, 1))
        self.model.decode_head.classifier = new_classifier
        self.model.to(self.gpu_id)
        for name, param in self.model.named_parameters():
            if 'classifier' not in name:
                param.requires_grad = False  # Freeze all parameters
        for param in self.model.decode_head.classifier.parameters():
            param.requires_grad = True  # unfreeze the new classifier
        trans_dim = args.bottleneck_dim
        self.transformer = MultiHeadAttentionOne(
            args.heads, trans_dim, trans_dim, trans_dim, dropout=args.dropout
        ).to(self.gpu_id)
        self.optimizer = get_optimizer([self.transformer,self.model.decode_head.classifier], args.optimizer, args.lr, args.weight_decay)
        iters_per_epoch = len(self.trainloader.dataset) // args.batch_size
        self.scheduler = get_scheduler(args.lr_scheduler, self.optimizer, args.epochs * iters_per_epoch, args.lr_power,
                                iters_per_epoch * args.lr_warmup, args.lr_warmup_ratio)
        self.infonce = InfoNCE(temperature=0.1, negative_mode='unpaired')
        self.confmat = utils.ConfusionMatrix(args.num_classes)

        if args.DDP:
            self.model = DDP(self.model, device_ids=[self.gpu_id])
            self.transformer = DDP(self.transformer, device_ids=[self.gpu_id])
    
    def _load_snapshot(self, snapshot_path):
        loc = f"cuda:{self.gpu_id}"
        snapshot = torch.load(snapshot_path, map_location=loc)
        self.model.load_state_dict(snapshot["MODEL_STATE"])
        self.epochs_run = snapshot["EPOCHS_RUN"]
        print(f"Resuming training from snapshot at Epoch {self.epochs_run}")

    def train_epoch(self, epoch):
        self.model.train()
        self.transformer.train()
        if self.args.DDP:
            self.trainloader.sampler.set_epoch(epoch)
        pbar = tqdm(self.trainloader)
        # prs = Presentation()
        # blank_slide_layout = prs.slide_layouts[6]
        # left = top = Inches(1)
        for i, batch in enumerate(pbar):
            q_img, q_label, spp_imgs, s_label, _, _, _ = batch
            q_img = q_img.to(self.gpu_id)
            q_label = q_label.to(self.gpu_id)
            spp_imgs = spp_imgs.to(self.gpu_id)
            s_label = s_label.to(self.gpu_id)

            # Phase 1: Train the binary classifier on support samples
            spp_imgs_reshape = spp_imgs.squeeze(1)  # [n_shots, 3, img_size, img_size]
            s_label_reshape = s_label.squeeze(1).long() # [n_shots, img_size, img_size]
            # No, use the classifier in Segformer

            optimizer = SGD(self.model.decode_head.classifier.parameters(), lr=args.lr)
            # Dynamic class weights
            s_label_arr = s_label.cpu().numpy().copy()  # [n_task, n_shots, img_size, img_size]
            back_pix = np.where(s_label_arr == 0)
            target_pix = np.where(s_label_arr == 1)

            if len(back_pix[0]) == 0 or len(target_pix[0]) == 0:
                continue  # skip bad support set
            criterion = nn.CrossEntropyLoss(
                weight=torch.tensor([1.0, len(back_pix[0]) / len(target_pix[0])]).to(self.gpu_id),
                ignore_index=255
            )
            s_outputs, s_hid = self.model(spp_imgs_reshape)  # [n_task, c, h, w]
            s_label_downsampled = F.max_pool2d(s_label.float(), 4)[0][0]  # yes, correct, continue [1,1,120,120]
            for index in range(args.adapt_iter):
                optimizer.zero_grad()
                spp_logits = self.model.decode_head.classifier(s_hid)
                support_output = F.interpolate(
                    spp_logits, size=s_label.size()[2:],
                    mode='bilinear', align_corners=True
                )
                s_loss = criterion(support_output, s_label_reshape)
                info_loss = None

                ##### 1.1 find FP, FN pixels, topk pos, topk neg pixels, done
                # inferred_spp_logits = spp_logits[0].argmax(0)
                # query_fn_coords = torch.nonzero((inferred_spp_logits != s_label_downsampled)&(s_label_downsampled == 1))
                # query_fp_coords = torch.nonzero((inferred_spp_logits != s_label_downsampled)&(s_label_downsampled == 0))
                # pos_coords = torch.nonzero((inferred_spp_logits == s_label_downsampled)&(s_label_downsampled == 1))
                # neg_coords = torch.nonzero((inferred_spp_logits == s_label_downsampled)&(s_label_downsampled == 0))
                # if query_fn_coords.size(0) < pos_coords.size(0):
                #     topk_1_indices = np.random.choice(pos_coords.size()[0], query_fn_coords.size(0), replace=False)
                # else:
                #     topk_1_indices = None
                # if query_fp_coords.size(0) < neg_coords.size(0):
                #     topk_0_indices = np.random.choice(neg_coords.size()[0], query_fp_coords.size(0), replace=False)
                # else:
                #     topk_0_indices = None
                # # 1.2 add infoNCE loss on q, p, n
                # if topk_1_indices is not None and topk_0_indices is not None and len(topk_1_indices) > 1 and len(topk_0_indices) > 1:
                #     topk_1_pixels = torch.permute(spp_logits[...,pos_coords[topk_1_indices][...,0],pos_coords[topk_1_indices][...,1]].squeeze(), (1,0))  # p
                #     topk_0_pixels = torch.permute(spp_logits[...,neg_coords[topk_0_indices][...,0],neg_coords[topk_0_indices][...,1]].squeeze(), (1,0))  # n
                #     q_fn_pixels = torch.permute(spp_logits[...,query_fn_coords[...,0],query_fn_coords[...,1]].squeeze(),(1,0))  # q_fn
                #     info_loss_fn = self.infonce(q_fn_pixels, topk_1_pixels, topk_0_pixels)
                #     q_fp_pixels = torch.permute(spp_logits[...,query_fp_coords[...,0],query_fp_coords[...,1]].squeeze(), (1,0))  # q_fp
                #     info_loss_fp = self.infonce(q_fp_pixels, topk_0_pixels, topk_1_pixels)
                #     info_loss = info_loss_fn + info_loss_fp
                
                # if info_loss:
                #     s_loss = s_loss + info_loss
    
                s_loss.backward()
                optimizer.step()
            
            # slide = prs.slides.add_slide(blank_slide_layout)
            # # 1. original image, 2. label mask, 3. inferred mask
            # ori_img = plt.imread(s_image_path_list[0][0])
            # fig, axs = plt.subplots(1,3,figsize=(10,3))
            # modified_label = s_label_downsampled[0][0].detach().cpu().numpy()
            # modified_label[modified_label == 255] = 0
            # axs[0].imshow(ori_img)
            # axs[0].axis('off')
            # axs[1].imshow(modified_label, cmap='viridis', alpha=0.4)
            # axs[1].axis('off')
            # axs[2].imshow(spp_logits[0].argmax(0).detach().cpu().numpy(), cmap='viridis', alpha=0.4)
            # axs[2].axis('off')
            # fig_filename = f'output/{i}_support.png'
            # plt.savefig(fig_filename)
            # plt.close()
            # pic = slide.shapes.add_picture(fig_filename, left, top)
            # if i > 10:
            #     break
            # Phase 2: Train the transformer to update the classifier's weights
            # Inputs of the transformer: weights of classifier trained on support sets, features of the query sample.
            # Dynamic class weights used for query image only during training
            q_label_arr = q_label.cpu().numpy().copy()  # [n_task, img_size, img_size]
            q_back_pix = np.where(q_label_arr == 0)
            q_target_pix = np.where(q_label_arr == 1)

            criterion = nn.CrossEntropyLoss(
                weight=torch.tensor([1.0, len(q_back_pix[0]) / (len(q_target_pix[0]) + 1e-12)]).to(self.gpu_id),
                ignore_index=255
            )

            # self.model.eval()
            # with torch.no_grad():
            _, q_hid = self.model(q_img)  # [n_task, c, h, w]
            q_hid = F.normalize(q_hid, dim=1)

            # Weights of the classifier.
            weights_cls = self.model.decode_head.classifier.weight.data  # [2,768,1,1]

            weights_cls_reshape = weights_cls.squeeze().unsqueeze(0).expand(
                args.batch_size, 2, weights_cls.shape[1]
            )  # [n_task, 2, c]

            # Update the classifier's weights with transformer
            updated_weights_cls = self.transformer(weights_cls_reshape, q_hid, q_hid)  # [n_task, 2, c]

            f_q_reshape = q_hid.view(args.batch_size, args.bottleneck_dim, -1)  # [n_task, c, hw]

            pred_q = torch.matmul(updated_weights_cls, f_q_reshape).view(
                args.batch_size, 2, q_hid.shape[-2], q_hid.shape[-1]
            )  # # [n_task, 2, h, w]

            pred_q = F.interpolate(
                pred_q, size=q_label.shape[1:],
                mode='bilinear', align_corners=True
            )

            q_loss = criterion(pred_q, q_label.long())

            self.optimizer.zero_grad()
            q_loss.backward()
            self.optimizer.step()

            pbar.set_description(f'e: {epoch}/{self.args.epochs}; iter: {i}; s_loss:{s_loss:.3f}, q_loss: {q_loss:.3f}')

        # prs.save(f'output/test.pptx')
        torch.cuda.empty_cache()

    def _save_snapshot(self, epoch):
        if self.args.DDP:
            snapshot = {
                "model": self.model.module.state_dict(),
                "transformer": self.transformer.module.state_dict(),
                "args": self.args,
            }
        else:
            snapshot = {
                "model": self.model.state_dict(),
                "transformer": self.transformer.state_dict(),
                "args": args,
            }
        save_path = f'{self.args.snapshot_path}_{epoch}.pt'
        torch.save(snapshot, save_path)
        print(f"Epoch {epoch} | Training snapshot saved at {save_path}")

    def train(self,):
        for epoch in range(self.args.epochs):
            self.train_epoch(epoch)
            if epoch % self.args.save_every == 0 and self.args.snapshot_path:
                self._save_snapshot(epoch)
            if epoch % self.args.eval_interval == 0:
                confmat = self.eval(epoch)
                val_info = str(confmat)
                print(val_info)

        confmat = self.eval(self.args.epochs)
        val_info = str(confmat)
        print(val_info)

    def eval(self, epoch):
        print(self.args.train_scenes, self.args.val_scenes)
        self.confmat.reset()
        self.model.train()
        self.transformer.eval()
        classes = [0,1]
        pbar = tqdm(self.valloader)
        for i, batch in enumerate(pbar):
            q_img, q_label, spp_imgs, s_label, subcls, spprt_oris, qry_oris = batch
            q_img = q_img.to(self.gpu_id)
            q_label = q_label.to(self.gpu_id)
            spp_imgs = spp_imgs.to(self.gpu_id)
            s_label = s_label.to(self.gpu_id)

            spp_imgs_reshape = spp_imgs.squeeze(1)  # [n_shots, 3, img_size, img_size]
            s_label_reshape = s_label.squeeze(1).long() # [n_shots, img_size, img_size]
            
            # Phase 1: Train a new binary classifier on support samples.
            binary_classifier = self.model.decode_head.classifier
            optimizer = SGD(binary_classifier.parameters(), lr=args.lr)
            # Dynamic class weights
            s_label_arr = s_label.cpu().numpy().copy()  # [n_task, n_shots, img_size, img_size]
            back_pix = np.where(s_label_arr == 0)
            target_pix = np.where(s_label_arr == 1)

            if len(back_pix[0]) == 0 or len(target_pix[0]) == 0:
                continue  # skip bad support set
            criterion = nn.CrossEntropyLoss(
                weight=torch.tensor([1.0, len(back_pix[0]) / len(target_pix[0])]).to(self.gpu_id),
                ignore_index=255
            )
            _, s_hid = self.model(spp_imgs_reshape)  # [n_task, c, h, w]

            for index in range(args.adapt_iter):
                optimizer.zero_grad()
                spp_logits = binary_classifier(s_hid)
                output_support = F.interpolate(
                    spp_logits, size=s_label.size()[2:],
                    mode='bilinear', align_corners=True
                )
                s_loss = criterion(output_support, s_label_reshape)
                s_loss.backward()
                optimizer.step()
            
            # Phase 2: Update classifier's weights with old weights and query features.
            _, q_hid = self.model(q_img)  # [n_task, c, h, w]
            q_hid = F.normalize(q_hid, dim=1)

            # Weights of the classifier.
            weights_cls = binary_classifier.weight.data  # [2,768,1,1]

            weights_cls_reshape = weights_cls.squeeze().unsqueeze(0).expand(
                q_hid.shape[0], 2, weights_cls.shape[1]
            )  # [n_task, 2, c]
            updated_weights_cls = self.transformer(weights_cls_reshape, q_hid, q_hid)

            # Build a temporary new classifier for prediction
            pseudo_cls = nn.Conv2d(
                args.bottleneck_dim, args.num_classes, kernel_size=1, bias=False
            ).to(self.gpu_id)

            pseudo_cls.weight.data = torch.as_tensor(
                updated_weights_cls.squeeze(0).unsqueeze(2).unsqueeze(3)
            )
            logits = pseudo_cls(q_hid)
            upsampled_logits = nn.functional.interpolate(logits, size=q_label.shape[-2:], mode="bilinear", align_corners=False)
            predicted = upsampled_logits.argmax(dim=1)
            self.confmat.update(q_label.flatten(), predicted.flatten())
            pbar.set_description(f'Epoch: {epoch}, eval iter: {i}')

        self.confmat.reduce_from_all_processes()
        self.wandb_log(self.confmat, epoch, self.args)
        return self.confmat
    
    def wandb_log(self, confmat, epoch, args):
        acc_global, acc, iu = confmat.compute()
        miou = iu.mean().item() * 100
        acc_global = acc_global.item() * 100
        acc = [x.item() for x in acc * 100]
        iu = [x.item() for x in iu * 100]
        log_dict = {'epoch': epoch, 'miou': miou, 'acc_global': acc_global,}
        log_dict['train_scenes'] = ';'.join(args.train_scenes)
        log_dict['val_scenes'] = ';'.join(args.val_scenes)
        for i,v in enumerate(acc):
            log_dict[f'acc_{i}'] = v
        
        for i,v in enumerate(iu):
            log_dict[f'iu_{i}'] = v
        
        wandb.log(log_dict)


def main(args):
    if args.DDP:
        ddp_setup()
    trainer = Trainer(args)
    trainer.train()
    if args.DDP:
        destroy_process_group()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        'Pytorch SegFormer Models training and evaluation script', parents=[get_argparser()])
    args = parser.parse_args()
    wandb.init(
        project=f'FSS_logits_align',
        config=args
    )
    for i, scene in enumerate(args.scenes):
        others = args.scenes[:i] + args.scenes[i+1:]  # yes
        args.train_scenes = [scene]
        args.val_scenes = others
        main(args)
