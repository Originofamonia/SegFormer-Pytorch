"""
https://huggingface.co/nvidia/segformer-b3-finetuned-cityscapes-1024-1024
https://huggingface.co/blog/fine-tune-segformer
torchrun --standalone --nproc_per_node=gpu SegFormer/segformer_trav_ddp.py
Draw
"""
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
import torch
import json
from huggingface_hub import hf_hub_download
import argparse
import numpy as np
import matplotlib.pyplot as plt
from transformers import SegformerForSemanticSegmentation
import multiprocessing as mp
from tqdm import tqdm
from torch.utils.data import DataLoader
from torch.nn import functional as F
import torch.nn as nn
from torch.optim import SGD
from torch.nn.parallel import DistributedDataParallel as DDP
from torch.utils.data import DistributedSampler, RandomSampler
from torch.distributed import init_process_group, destroy_process_group

from local_datasets.cityscapes import trav_train_loader, trav_val_loader
from models.transformer import MultiHeadAttentionOne
from models.infonce import InfoNCE
from utils.schedulers import get_scheduler, create_lr_scheduler
from utils.optimizers import get_optimizer
# from utils.utils import fix_seeds, setup_cudnn, cleanup_ddp, setup_ddp
import utils.distributed_utils as utils

# os.environ['OMP_NUM_THREADS'] = f"{int(os.cpu_count() // torch.cuda.device_count())}"  # no, slower

def ddp_setup():
    init_process_group(backend="nccl")
    torch.cuda.set_device(int(os.environ["LOCAL_RANK"]))


def get_argparser():
    parser = argparse.ArgumentParser('Pytorch SegFormer Models training and evaluation script', add_help=False)

    # Datset Options
    parser.add_argument("--data_root", type=str, default='/mnt/hdd/segmentation_indoor_images',help="path to Dataset")
    parser.add_argument("--scenes", type=list, default=['elb', 'erb', 'uc', 'woh'],
                        choices=['elb', 'erb', 'uc', 'nh', 'woh'], help='Name of dataset')
    parser.add_argument("--image_size", type=list, default=480, help="input size, [480, 640]")
    parser.add_argument("--ignore_label", type=int, default=255, help="path to Dataset")
    parser.add_argument("--dataset", type=str, default='trav',choices=['cityscapes', 'trav'], help='Name of dataset')
    parser.add_argument("--num_classes", type=int, default=2, help="num classes (default: None)")
    parser.add_argument("--pin_mem", type=bool, default=True, help="Dataloader ping_memory")
    parser.add_argument("--batch_size", type=int, default=1,help='batch size (zeus:10, poseidon:4)')
    parser.add_argument("--val_batch_size", type=int, default=4,help='batch size for validation')
    parser.add_argument("--scale_min", type=float, default=0.5)
    parser.add_argument("--scale_max", type=float, default=2.0)
    parser.add_argument("--rot_min", type=float, default=-10)
    parser.add_argument("--rot_max", type=float, default=10)
    parser.add_argument("--mean", type=list, default=[0.5174, 0.4857, 0.5054])
    parser.add_argument("--std", type=list, default=[0.2726, 0.2778, 0.2861])
    parser.add_argument("--augmentations", type=list, default=[])
    parser.add_argument("--train_split", type=int, default=0)
    parser.add_argument("--test_split", type=str, default='default')
    parser.add_argument("--random_shot", type=bool, default=False)
    parser.add_argument("--shot", type=int, default=1)

    # Model Options
    # parser.add_argument("--model", type=str, default='make_SegFormerB1', help='model name')
    parser.add_argument("--bottleneck_dim", type=int, default=768)
    parser.add_argument("--heads", type=int, default=1)
    parser.add_argument("--dropout", type=float, default=0.5)

    # Train Options
    parser.add_argument("--amp", type=bool, default=False, help='auto mixture precision, do not use') # There may be some problems when loading weights, such as: ComplexFloat
    parser.add_argument("--epochs", type=int, default=20, help='total training epochs')
    parser.add_argument("--device", type=str, default='cuda:0', help='device (cuda:0 or cpu)')
    parser.add_argument("--workers", type=int, default=4,
                        help='workers, set it equal 0 when run programs in win platform')
    parser.add_argument("--DDP", type=bool, default=False)
    parser.add_argument("--train_print_freq", type=int, default=50)
    parser.add_argument("--val_print_freq", type=int, default=50)
    parser.add_argument("--adapt_iter", type=int, default=200)

    # Loss Options
    parser.add_argument("--loss_fn_name", type=str, default='CrossEntropy', choices=['OhemCrossEntropy', 'CrossEntropy'])

    # Optimizer & LR-scheduler Options
    parser.add_argument("--optimizer", type=str, default='adamw')
    parser.add_argument("--lr", type=float, default=2.5e-4,
                        help="learning rate (default: 1e-3)")
    parser.add_argument("--weight_decay", type=float, default=1e-4,
                        help='weight decay (default: 1e-5)')

    parser.add_argument("--lr_scheduler", type=str, default='WarmupPolyLR')
    parser.add_argument("--lr_power", type=float, default=0.9)
    parser.add_argument("--lr_warmup", type=int, default=10)
    parser.add_argument("--lr_warmup_ratio", type=float, default=0.1)

    # save checkpoints
    parser.add_argument("--save_weights_dir", default='./save_weights', type=str,
                        help="restore from checkpoint")
    parser.add_argument("--resume", type=bool, default=False)
    parser.add_argument("--gpu_id", type=int, default=0, help="GPU ID")
    parser.add_argument("--save_dir", type=str, default='./', help='SummaryWriter save dir')
    parser.add_argument("--eval_interval", type=int, default=2, help="evaluation interval")
    parser.add_argument("--save_every", type=int, default=2, help="evaluation interval")
    parser.add_argument("--load_pretrained", type=bool, default=True)
    parser.add_argument("--pretrained_path", type=str, default='nvidia/segformer-b3-finetuned-cityscapes-1024-1024')
    parser.add_argument("--snapshot_path", type=str, default='save_weights/trav_fs_infonce_10.pt')
    parser.add_argument("--column", type=str, default='fs_10')
    return parser

class Trainer:
    def __init__(self, args) -> None:
        self.args = args
        # init train objects
        if args.DDP:
            self.gpu_id = int(os.environ["LOCAL_RANK"])
        else:
            self.gpu_id = args.device

        repo_id = "huggingface/label-files"
        id2label_filename = "cityscapes-id2label.json"
        id2label = json.load(open(hf_hub_download(repo_id, id2label_filename, repo_type="dataset"), "r"))
        id2label = {int(k): v for k, v in id2label.items()}
        label2id = {v: k for k, v in id2label.items()}

        self.trainloader, _ = trav_train_loader(args)
        self.valloader, _ = trav_val_loader(args)

        self.model = SegformerForSemanticSegmentation.from_pretrained(args.pretrained_path,
            id2label=id2label, label2id=label2id)
        new_classifier = nn.Conv2d(768, 2, kernel_size=(1, 1), stride=(1, 1))
        self.model.decode_head.classifier = new_classifier
        self.model.to(self.gpu_id)
        for name, param in self.model.named_parameters():
            if 'classifier' not in name:
                param.requires_grad = False  # Freeze all parameters
        for param in self.model.decode_head.classifier.parameters():
            param.requires_grad = True  # unfreeze the new classifier
        trans_dim = args.bottleneck_dim
        self.transformer = MultiHeadAttentionOne(
            args.heads, trans_dim, trans_dim, trans_dim, dropout=args.dropout
        ).to(self.gpu_id)
        self.optimizer = get_optimizer([self.transformer,self.model.decode_head.classifier], args.optimizer, args.lr, args.weight_decay)
        iters_per_epoch = len(self.trainloader.dataset) // args.batch_size
        self.scheduler = get_scheduler(args.lr_scheduler, self.optimizer, args.epochs * iters_per_epoch, args.lr_power,
                                iters_per_epoch * args.lr_warmup, args.lr_warmup_ratio)
        self.infonce = InfoNCE(temperature=0.1, negative_mode='unpaired')
        self.confmat = utils.ConfusionMatrix(args.num_classes)
        self._load_snapshot(args.snapshot_path)

        if args.DDP:
            self.model = DDP(self.model, device_ids=[self.gpu_id])
            self.transformer = DDP(self.transformer, device_ids=[self.gpu_id])
    
    def _load_snapshot(self, snapshot_path):
        loc = f"{self.gpu_id}"
        snapshot = torch.load(snapshot_path, map_location=loc)
        self.model.load_state_dict(snapshot['model'])
        self.transformer.load_state_dict(snapshot['transformer'])
        # self.epochs_run = snapshot["EPOCHS_RUN"]
        print(f"Load snapshot at {snapshot_path}")

    def create_directory_if_not_exists(self, path):
        # Check if the directory exists
        if not os.path.exists(path):
            # If not, create it (including any necessary parent directories)
            os.makedirs(path)
            print(f"Directory '{path}' created.")
        else:
            print(f"Directory '{path}' already exists.")

    def infer(self):
        """
        save inferred masks from finetuned models
        """
        self.confmat.reset()
        self.model.train()
        self.transformer.eval()
        pbar = tqdm(self.valloader)
        self.create_directory_if_not_exists(f'output/{self.args.column}')
        for i, batch in enumerate(pbar):
            q_img, q_label, spp_imgs, s_label, subcls, spp_oris, q_oris = batch
            q_img = q_img.to(self.gpu_id)
            q_label = q_label.to(self.gpu_id)
            spp_imgs = spp_imgs.to(self.gpu_id)
            s_label = s_label.to(self.gpu_id)

            spp_imgs_reshape = spp_imgs.squeeze(1)  # [n_shots, 3, img_size, img_size]
            s_label_reshape = s_label.squeeze(1).long() # [n_shots, img_size, img_size]
            
            # Phase 1: Train a new binary classifier on support samples.
            binary_classifier = self.model.decode_head.classifier
            optimizer = SGD(binary_classifier.parameters(), lr=self.args.lr)
            # Dynamic class weights
            s_label_arr = s_label.cpu().numpy().copy()  # [n_task, n_shots, img_size, img_size]
            back_pix = np.where(s_label_arr == 0)
            target_pix = np.where(s_label_arr == 1)

            if len(back_pix[0]) == 0 or len(target_pix[0]) == 0:
                continue  # skip bad support set
            criterion = nn.CrossEntropyLoss(
                weight=torch.tensor([1.0, len(back_pix[0]) / len(target_pix[0])]).to(self.gpu_id),
                ignore_index=255
            )
            _, s_hid = self.model(spp_imgs_reshape)  # [n_task, c, h, w]

            for index in range(self.args.adapt_iter):
                optimizer.zero_grad()
                spp_logits = binary_classifier(s_hid)
                output_support = F.interpolate(
                    spp_logits, size=s_label.size()[2:],
                    mode='bilinear', align_corners=True
                )
                s_loss = criterion(output_support, s_label_reshape)
                s_loss.backward()
                optimizer.step()
            
            # Phase 2: Update classifier's weights with old weights and query features.
            _, q_hid = self.model(q_img)  # [n_task, c, h, w]
            q_hid = F.normalize(q_hid, dim=1)

            # Weights of the classifier.
            weights_cls = binary_classifier.weight.data  # [2,768,1,1]

            weights_cls_reshape = weights_cls.squeeze().unsqueeze(0).expand(
                q_hid.shape[0], 2, weights_cls.shape[1]
            )  # [n_task, 2, c]
            updated_weights_cls = self.transformer(weights_cls_reshape, q_hid, q_hid)

            # Build a temporary new classifier for prediction
            pseudo_cls = nn.Conv2d(
                self.args.bottleneck_dim, self.args.num_classes, kernel_size=1, bias=False
            ).to(self.gpu_id)

            pseudo_cls.weight.data = torch.as_tensor(
                updated_weights_cls.squeeze(0).unsqueeze(2).unsqueeze(3)
            )
            q_logits = pseudo_cls(q_hid)
            upsampled_logits = nn.functional.interpolate(q_logits, size=q_label.shape[-2:], mode="bilinear", align_corners=False)
            q_pred = upsampled_logits.argmax(dim=1)[0].detach().cpu().numpy()
            # save q_pred to individual images
            q_pred_filename = q_oris[0][0].split('/')[-1].strip('.jpg')
            np.save(f'output/{self.args.column}/{q_pred_filename}_{self.args.column}', q_pred)
            pbar.set_description(f'Iter: {i}')

        self.confmat.reduce_from_all_processes()
        return self.confmat
    
    def draw_pptx(self,):
        """
        each ppt page: 0,4,6,10,16,18
        """
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]
        left = top = Inches(0.1)
        top_2 = Inches(6)
        top_3 = Inches(8)
        alpha = 0.6
        width = Inches(14.0)
        height = Inches(1.2)
        pbar = tqdm(self.valloader)
        for i, batch in enumerate(pbar):
            _, _, _, _, _, spp_oris, q_oris = batch
            q_pred_filename = q_oris[0][0].split('/')[-1].strip('.jpg')
            s_filename = spp_oris[0][0][0]
            fs_0_filename = f'output/fs_0/{q_pred_filename}_fs_0.npy'
            fs_4_filename = f'output/fs_4/{q_pred_filename}_fs_4.npy'
            fs_6_filename = f'output/fs_6/{q_pred_filename}_fs_6.npy'
            fs_10_filename = f'output/fs_10/{q_pred_filename}_fs_10.npy'
            fs_16_filename = f'output/fs_16/{q_pred_filename}_fs_16.npy'
            fs_18_filename = f'output/fs_18/{q_pred_filename}_fs_18.npy'
            slide = prs.slides.add_slide(blank_slide_layout)
            fig, axs = plt.subplots(2, 5, figsize=(14, 6))  # w,h
            q_img = plt.imread(q_oris[0][0])
            q_target = np.load(q_oris[1][0])
            s_img = plt.imread(s_filename)
            s_target = np.load(spp_oris[1][0][0])
            fs_0 = np.load(fs_0_filename)
            fs_4 = np.load(fs_4_filename)
            fs_6 = np.load(fs_6_filename)
            fs_10 = np.load(fs_10_filename)
            fs_16 = np.load(fs_16_filename)
            fs_18 = np.load(fs_18_filename)
            
            axs[0,0].imshow(s_img)
            axs[0,0].set_title(f's_img')
            axs[0,0].axis('off')
            
            axs[1,0].imshow(s_img)
            axs[1,0].imshow(s_target, cmap='viridis', alpha=alpha)
            axs[1,0].set_title(f's_target')
            axs[1,0].axis('off')

            axs[0,1].imshow(q_img)
            axs[0,1].set_title(f'q_img')
            axs[0,1].axis('off')
            
            axs[1,1].imshow(q_img)
            axs[1,1].imshow(q_target, cmap='viridis', alpha=alpha)
            axs[1,1].set_title(f'q_target')
            axs[1,1].axis('off')

            axs[0,2].imshow(q_img)
            axs[0,2].imshow(fs_0, cmap='viridis', alpha=alpha)
            axs[0,2].set_title(f'fs_0')
            axs[0,2].axis('off')

            axs[0,3].imshow(q_img)
            axs[0,3].imshow(fs_4, cmap='viridis', alpha=alpha)
            axs[0,3].set_title(f'fs_4')
            axs[0,3].axis('off')

            axs[0,4].imshow(q_img)
            axs[0,4].imshow(fs_6, cmap='viridis', alpha=alpha)
            axs[0,4].set_title(f'fs_6')
            axs[0,4].axis('off')

            axs[1,2].imshow(q_img)
            axs[1,2].imshow(fs_10, cmap='viridis', alpha=alpha)
            axs[1,2].set_title(f'fs_10')
            axs[1,2].axis('off')

            axs[1,3].imshow(q_img)
            axs[1,3].imshow(fs_16, cmap='viridis', alpha=alpha)
            axs[1,3].set_title(f'fs_16')
            axs[1,3].axis('off')

            axs[1,4].imshow(q_img)
            axs[1,4].imshow(fs_18, cmap='viridis', alpha=alpha)
            axs[1,4].set_title(f'fs_18')
            axs[1,4].axis('off')
            plt.subplots_adjust(hspace=0.01, wspace=0.01)
            img_filename = f'output/pptx/{q_pred_filename}.png'
            fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
            plt.close()
            pic = slide.shapes.add_picture(img_filename, left, top)
            text_box = slide.shapes.add_textbox(left, top_2, width, height)

            # Get the text frame within the text box
            tf = text_box.text_frame

            # Add a paragraph to the text frame
            p = tf.add_paragraph()
            p.text = f'q: {q_oris[0][0]};\n s: {s_filename}'
        
        prs.save(f'output/different_models.pptx')

# Use Computer Modern font
plt.rcParams['font.family'] = 'serif'
plt.rcParams['font.serif'] = ['Computer Modern Roman']
plt.rcParams['text.usetex'] = True

def draw_qual():
    pos_files = ['/mnt/hdd/segmentation_indoor_images/elb/challenging/images/1664300444727277987.jpg',
                 '/mnt/hdd/segmentation_indoor_images/erb/challenging/images/1664303453189311735.jpg',
                 '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661556043781475060.jpg',
                 '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661555947496401064.jpg',
                 '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661556434679290525.jpg',
                 '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661555874275943510.jpg']
    neg_files = ['/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661556016257148735.jpg',
                 '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661555927693561069.jpg',]
    font_size = 6
    alpha = 0.6
    fig, axs = plt.subplots(6, 4, figsize=(6.5, 7.5))  # w, h
    for i, img in enumerate(pos_files):
        ori_img = plt.imread(img)
        q_pred_filename = img.split('/')[-1].strip('.jpg')
        target_filename = img.replace('/images', '/labels', 1)
        target_filename = target_filename.replace('.jpg', '.npy')
        fs_0_filename = f'output/fs_0/{q_pred_filename}_fs_0.npy'
        fs_18_filename = f'output/fs_18/{q_pred_filename}_fs_18.npy'
        target = np.load(target_filename)
        fs_0 = np.load(fs_0_filename)
        fs_18 = np.load(fs_18_filename)
        axs[i,0].imshow(ori_img)
        axs[i,0].axis('off')
        
        axs[i,1].imshow(ori_img)
        axs[i,1].imshow(target, cmap='viridis', alpha=alpha)
        axs[i,1].axis('off')

        axs[i,2].imshow(ori_img)
        axs[i,2].imshow(fs_0, cmap='viridis', alpha=alpha)
        axs[i,2].axis('off')

        axs[i,3].imshow(ori_img)
        axs[i,3].imshow(fs_18, cmap='viridis', alpha=alpha)
        axs[i,3].axis('off')
    column_text = ["Query", 'Target', 'FSS', 'FSS+CPC (ours)']
    for j in range(4):
        axs[0, j].text(0.5, 1.1, f'{column_text[j]}', ha='center', va='center', transform=axs[0, j].transAxes, fontsize=font_size)
    
    # Adjust layout for better spacing
    plt.subplots_adjust(hspace=0.05, wspace=0.05)
    img_filename = f'output/qualitative/positive.pdf'
    fig.savefig(img_filename, bbox_inches='tight', pad_inches=0.05)
    plt.close()

    fig, axs = plt.subplots(2, 4, figsize=(6, 2.3))  # w, h
    for i, img in enumerate(neg_files):
        ori_img = plt.imread(img)
        q_pred_filename = img.split('/')[-1].strip('.jpg')
        target_filename = img.replace('/images', '/labels', 1)
        target_filename = target_filename.replace('.jpg', '.npy')
        fs_0_filename = f'output/fs_0/{q_pred_filename}_fs_0.npy'
        fs_18_filename = f'output/fs_18/{q_pred_filename}_fs_18.npy'
        target = np.load(target_filename)
        fs_0 = np.load(fs_0_filename)
        fs_18 = np.load(fs_18_filename)
        axs[i,0].imshow(ori_img)
        axs[i,0].axis('off')

        axs[i,1].imshow(ori_img)
        axs[i,1].imshow(target, cmap='viridis', alpha=alpha)
        axs[i,1].axis('off')

        axs[i,2].imshow(ori_img)
        axs[i,2].imshow(fs_0, cmap='viridis', alpha=alpha)
        axs[i,2].axis('off')

        axs[i,3].imshow(ori_img)
        axs[i,3].imshow(fs_18, cmap='viridis', alpha=alpha)
        axs[i,3].axis('off')
    
    for j in range(4):
        axs[0, j].text(0.5, 1.1, f'{column_text[j]}', ha='center', va='center', transform=axs[0, j].transAxes, fontsize=font_size)
    
    # Adjust layout for better spacing
    plt.subplots_adjust(hspace=0.05, wspace=0.05)
    img_filename = f'output/qualitative/negative.pdf'
    fig.savefig(img_filename, bbox_inches='tight', pad_inches=0.05)
    plt.close()


def draw_arch_cpc():
    q_img = '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661556141830365889.jpg'
    q_pred_filename = q_img.split('/')[-1].strip('.jpg')
    q_target_filename = q_img.replace('/images', '/labels', 1)
    q_target_filename = q_target_filename.replace('.jpg', '.npy')
    q_fs_0_filename = f'output/fs_18/{q_pred_filename}_fs_18.npy'
    q_target = np.load(q_target_filename)
    q_pred = np.load(q_fs_0_filename)
    q_diff = q_target - q_pred
    s_img = '/mnt/hdd/segmentation_indoor_images/uc/challenging/images/1661555948062233679.jpg'
    s_pred_filename = s_img.split('/')[-1].strip('.jpg')
    s_target_filename = s_img.replace('/images', '/labels', 1)
    s_target_filename = s_target_filename.replace('.jpg', '.npy')
    s_fs_0_filename = f'output/fs_18/{s_pred_filename}_fs_18.npy'
    s_target = np.load(s_target_filename)
    s_pred = np.load(s_fs_0_filename)
    s_diff = s_target - s_pred
    s_color = np.zeros_like(s_diff)
    # Assign different values based on TP, TN, FP, FN
    s_color[(s_target == 1) & (s_pred == 1)] = 3  # True Positive (TP)
    s_color[(s_target == 0) & (s_pred == 0)] = 0  # True Negative (TN)
    s_color[(s_target == 0) & (s_pred == 1)] = 1  # False Positive (FP)
    s_color[(s_target == 1) & (s_pred == 0)] = 2  # False Negative (FN)
    fig, ax = plt.subplots()
    im = ax.imshow(s_color, cmap='Accent', alpha=0.9)
    ax.axis('off')
    img_filename = f'output/qualitative/s_contrast_18.png'
    fig.savefig(img_filename, bbox_inches='tight', pad_inches=0.0)


def main(args):
    if args.DDP:
        ddp_setup()
    trainer = Trainer(args)
    # trainer.infer()
    trainer.draw_pptx()
    if args.DDP:
        destroy_process_group()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        'Pytorch SegFormer Models training and evaluation script', parents=[get_argparser()])
    args = parser.parse_args()
    main(args)
    # draw_qual()
    # draw_arch_cpc()
